"""
PPTXGenerator - Génère le one-pager Weekly Status PowerPoint.
Basé sur le template Project template TMM.pptx
"""
import copy
import io
import logging
from datetime import date, datetime
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx non disponible")

TEMPLATE_PATH = Path("templates/weekly_template.pptx")

KPI_HEX = {"G": "92D050", "Y": "FFFF00", "R": "FF0000"}

MONTHS_FR = ["", "Jan.", "Fév.", "Mars", "Avr.", "Mai", "Juin",
             "Juil.", "Août", "Sept.", "Oct.", "Nov.", "Déc."]
MONTHS_EN = ["", "Jan.", "Feb.", "Mar.", "Apr.", "May", "Jun.",
             "Jul.", "Aug.", "Sep.", "Oct.", "Nov.", "Dec."]


def _fmt_date(d: date, language: str = "fr") -> str:
    months = MONTHS_FR if language == "fr" else MONTHS_EN
    suffix = "st" if d.day == 1 else "nd" if d.day == 2 else "rd" if d.day == 3 else "th"
    if language == "fr":
        return f"{d.day} {months[d.month]} {d.year}"
    return f"{months[d.month]} {d.day}{suffix}, {d.year}"


def generate_weekly_pptx(
    project:    dict,
    tasks:      list[dict],
    categories: list[dict],
    kpis:       list[dict],
    milestones: list[dict],
    risks:      list[dict],
    weekly_note: dict | None,
    language:   str = "fr",
) -> bytes:
    """
    Génère le fichier PowerPoint Weekly Status.
    Retourne les bytes du fichier .pptx.
    """
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx non installé")

    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Template introuvable : {TEMPLATE_PATH}")

    prs  = Presentation(str(TEMPLATE_PATH))
    slide = prs.slides[0]
    today = date.today()

    def set_text(shape_name: str, text: str) -> None:
        for shape in slide.shapes:
            if shape.name == shape_name and shape.has_text_frame:
                tf = shape.text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            run.text = text
                            return

    # ── Titre et date ─────────────────────────────────────────────────────────
    for shape in slide.shapes:
        if shape.name == "Title 1" and shape.has_text_frame:
            tf = shape.text_frame
            paras = tf.paragraphs
            if len(paras) >= 2:
                for run in paras[1].runs:
                    run.text = f"       {project.get('name', '')}"

        if shape.name == "TextBox 12" and shape.has_text_frame:
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = f"              {_fmt_date(today, language)}"

        if shape.name == "TextBox 1" and shape.has_text_frame:
            go_live = project.get("go_live_date", "TBD")
            lbl = "Go-Live" if language == "en" else "Mise en prod."
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = f"{lbl} : {go_live}"

    # ── KPIs (Table 26) ───────────────────────────────────────────────────────
    kpi_map = {k["kpi_name"]: k for k in kpis}
    for shape in slide.shapes:
        if shape.name == "Table 26" and shape.has_table:
            table = shape.table
            kpi_order = ["COST", "SCOPE", "SCHEDULE", "RESOURCES", "RISK", "TRANSITION"]
            for row_idx, kpi_name in enumerate(kpi_order, start=1):
                kpi = kpi_map.get(kpi_name, {})
                prev_val = kpi.get("prev_value", "G")
                curr_val = kpi.get("curr_value", "G")
                for col_idx, val in [(1, prev_val), (2, curr_val)]:
                    cell = table.cell(row_idx, col_idx)
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = val
                    # Couleur de fond
                    hex_color = KPI_HEX.get(val, "92D050")
                    fill = cell.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor.from_string(hex_color)

    # ── Summary (Table 7) ────────────────────────────────────────────────────
    if weekly_note:
        for shape in slide.shapes:
            if shape.name == "Table 7" and shape.has_table:
                cell = shape.table.cell(1, 0)
                cell.text_frame.paragraphs[0].runs[0].text =                     weekly_note.get("summary", "")

    # ── Achievements (Table 25) ───────────────────────────────────────────────
    if weekly_note:
        for shape in slide.shapes:
            if shape.name == "Table 25" and shape.has_table:
                cell = shape.table.cell(1, 0)
                text = weekly_note.get("achievements", "")
                if cell.text_frame.paragraphs:
                    cell.text_frame.paragraphs[0].runs[0].text = text

    # ── Planned (Table 13) ────────────────────────────────────────────────────
    if weekly_note:
        for shape in slide.shapes:
            if shape.name == "Table 13" and shape.has_table:
                cell = shape.table.cell(1, 0)
                text = weekly_note.get("planned", "")
                if cell.text_frame.paragraphs:
                    cell.text_frame.paragraphs[0].runs[0].text = text

    # ── Watch Items (Table 15) ────────────────────────────────────────────────
    if weekly_note:
        for shape in slide.shapes:
            if shape.name == "Table 15" and shape.has_table:
                cell = shape.table.cell(1, 0)
                text = weekly_note.get("watch_items", "")
                if cell.text_frame.paragraphs:
                    cell.text_frame.paragraphs[0].runs[0].text = text

    # ── Milestones (Table 4) ──────────────────────────────────────────────────
    for shape in slide.shapes:
        if shape.name == "Table 4" and shape.has_table:
            cell_titles   = shape.table.cell(1, 0)
            cell_baseline = shape.table.cell(1, 1)
            cell_current  = shape.table.cell(1, 2)
            cell_status   = shape.table.cell(1, 3)
            titles    = "\n".join(m.get("title", "")         for m in milestones)
            baselines = "\n".join(m.get("baseline_date", "") for m in milestones)
            currents  = "\n".join(m.get("current_date", "")  for m in milestones)
            statuses  = "\n".join(m.get("status", "")        for m in milestones)
            for cell, text in [(cell_titles, titles), (cell_baseline, baselines),
                               (cell_current, currents), (cell_status, statuses)]:
                if cell.text_frame.paragraphs:
                    try:
                        cell.text_frame.paragraphs[0].runs[0].text = text
                    except IndexError:
                        pass

    # ── Risks (Table 20) ──────────────────────────────────────────────────────
    for shape in slide.shapes:
        if shape.name == "Table 20" and shape.has_table:
            cell_desc  = shape.table.cell(1, 0)
            cell_owner = shape.table.cell(1, 1)
            descs  = "\n".join(f"{r.get('risk_type','I')}:  {r.get('description','')}" for r in risks)
            owners = "\n".join(r.get("owner", "") for r in risks)
            try:
                cell_desc.text_frame.paragraphs[0].runs[0].text  = descs
                cell_owner.text_frame.paragraphs[0].runs[0].text = owners
            except IndexError:
                pass

    # Sauvegarder en mémoire
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
